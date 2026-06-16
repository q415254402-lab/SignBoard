"""PPT 转图片序列模块"""

import os
import uuid
import subprocess
import sys
from pathlib import Path

from PIL import Image


def convert_pptx_to_images(pptx_path: str, output_dir: str) -> list[str]:
    """
    将 PPTX 文件转为图片序列。
    返回图片文件名列表（相对路径）。

    优先使用 LibreOffice（无头转换），
    如果不可用则尝试 python-pptx + Pillow（效果有限）。
    """
    os.makedirs(output_dir, exist_ok=True)

    # 方法1：尝试 LibreOffice
    images = _convert_with_libreoffice(pptx_path, output_dir)
    if images:
        return images

    # 方法2：python-pptx fallback
    images = _convert_with_pptx(pptx_path, output_dir)
    return images


def _convert_with_libreoffice(pptx_path: str, output_dir: str) -> list[str]:
    """使用 LibreOffice 转为图片（最可靠）"""
    # 尝试找 LibreOffice
    libreoffice_paths = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "soffice",
        "libreoffice",
    ]

    soffice = None
    for p in libreoffice_paths:
        try:
            result = subprocess.run([p, "--version"], capture_output=True, timeout=5)
            if result.returncode == 0:
                soffice = p
                break
        except Exception:
            continue

    if not soffice:
        return []

    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "png",
             f"--outdir", output_dir, pptx_path],
            timeout=120, check=True, capture_output=True,
        )
    except subprocess.TimeoutExpired:
        return []
    except Exception:
        return []

    # 收集生成的图片（只取本次转换后生成的文件，避免污染）
    ppt_mtime = os.path.getmtime(pptx_path)
    images = sorted([
        f for f in os.listdir(output_dir)
        if f.lower().endswith(('.png', '.jpg', '.jpeg'))
        and os.path.getmtime(os.path.join(output_dir, f)) >= ppt_mtime
    ])

    return images


def _convert_with_pptx(pptx_path: str, output_dir: str) -> list[str]:
    """使用 python-pptx 提取图片（fallback，效果有限）"""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        return []

    try:
        prs = Presentation(pptx_path)
    except Exception:
        return []

    images = []

    for i, slide in enumerate(prs.slides):
        slide_width = prs.slide_width or Inches(13.333)
        slide_height = prs.slide_height or Inches(7.5)

        # 创建空白画布
        img = Image.new("RGB", (
            int(slide_width / 914400 * 96),  # EMU to pixels at 96 DPI
            int(slide_height / 914400 * 96),
        ), "white")

        for shape in slide.shapes:
            if shape.has_table:
                continue
            if shape.shape_type == 13:  # Picture
                try:
                    image_part = shape.image
                    image_bytes = image_part.blob
                    from io import BytesIO
                    pil_img = Image.open(BytesIO(image_bytes))
                    # 调整大小
                    left = int(shape.left / 914400 * 96) if shape.left else 0
                    top = int(shape.top / 914400 * 96) if shape.top else 0
                    w = int(shape.width / 914400 * 96) if shape.width else pil_img.width
                    h = int(shape.height / 914400 * 96) if shape.height else pil_img.height
                    pil_img = pil_img.resize((w, h), Image.LANCZOS)
                    img.paste(pil_img, (left, top))
                except Exception:
                    pass

        filename = f"pptx_slide_{i+1:03d}.png"
        filepath = os.path.join(output_dir, filename)
        img.save(filepath, quality=90)
        images.append(filename)

    return images


def convert_pptx_upload(pptx_path: str, upload_dir: str) -> tuple[str, list[str]]:
    """
    处理上传的 PPTX：转换为图片序列并存到 upload_dir。
    返回 (目录名, 图片文件名列表)。
    """
    ppt_id = uuid.uuid4().hex[:8]
    ppt_dir_name = f"ppt_{ppt_id}"
    ppt_output_dir = os.path.join(upload_dir, ppt_dir_name)
    images = convert_pptx_to_images(pptx_path, ppt_output_dir)
    return ppt_dir_name, images